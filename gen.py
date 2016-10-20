from pptx import Presentation
import datetime as dt

from ppt_utils.days import next_days

if __name__ == "__main__":
    # Choose your month
    month = dt.date(day=1, month=12, year=2016)

    # Find all the dates
    days = next_days(['Mon', 'Wed', 'Fri'], month)

    # Create a presentation and choose a layout
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    # Make slides with dates on them
    for day in days:
        slide = presentation.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = day.strftime('%a %b %d %Y')

    # Create a new presentation
    presentation.save('test.pptx')
